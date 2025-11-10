"""
Lightweight document processor for extracting text from various file formats.
Replaces docling with lighter alternatives to reduce Docker image size.
"""

import io
import json
import tempfile
import os

# PDF processing
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# Office document processing
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import docx2txt
except ImportError:
    docx2txt = None

from app.shared.logger.logger import console_log


class DocumentProcessor:
    """Lightweight document processor for text extraction from various formats."""

    def __init__(self):
        console_log("Initializing lightweight DocumentProcessor")

    def extract_text_from_pdf(self, file_data: bytes, source_name: str) -> str:
        """Extract text from PDF using PyMuPDF (fitz)."""
        if not fitz:
            return "PyMuPDF not available for PDF processing"

        try:
            # Validate that we have valid PDF data
            if not file_data or len(file_data) < 100:
                return "Invalid PDF data: file too small or empty"

            # Check if it starts with PDF header
            if not file_data.startswith(b"%PDF-"):
                return "Invalid PDF data: does not start with PDF header"

            console_log(f"Processing PDF with {len(file_data)} bytes for: {source_name}")

            # Try multiple approaches to open the PDF
            pdf_document = None

            # Approach 1: Try with stream parameter and explicit filetype
            try:
                pdf_document = fitz.open(stream=file_data, filetype="pdf")
                console_log("Successfully opened PDF using stream method")
            except Exception as stream_error:
                console_log(f"Stream method failed: {stream_error}, trying BytesIO method")

                # Approach 2: Try with BytesIO
                try:
                    import io

                    pdf_document = fitz.open(stream=io.BytesIO(file_data), filetype="pdf")
                    console_log("Successfully opened PDF using BytesIO method")
                except Exception as bytesio_error:
                    console_log(f"BytesIO method failed: {bytesio_error}, trying memory method")

                    # Approach 3: Try opening from memory without filetype
                    try:
                        pdf_document = fitz.open(stream=file_data)
                        console_log("Successfully opened PDF using memory method")
                    except Exception as memory_error:
                        console_log(f"Memory method failed: {memory_error}, trying temporary file method")

                        # Approach 4: Try with temporary file
                        try:
                            import tempfile

                            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                                temp_file.write(file_data)
                                temp_file.flush()
                                pdf_document = fitz.open(temp_file.name)
                                console_log("Successfully opened PDF using temporary file method")
                                # Clean up the temp file
                                import os

                                os.unlink(temp_file.name)
                        except Exception as temp_error:
                            console_log(f"Temporary file method failed: {temp_error}")
                            return f"Error opening PDF: All methods failed. Stream: {stream_error}, BytesIO: {bytesio_error}, Memory: {memory_error}, Temp: {temp_error}"

            if not pdf_document:
                return "Failed to open PDF document"

            text_content = ""
            page_count = pdf_document.page_count
            console_log(f"PDF has {page_count} pages")

            for page_num in range(page_count):
                try:
                    page = pdf_document[page_num]
                    # Use getattr to safely access the get_text method
                    page_text = getattr(page, "get_text", lambda: "")()
                    text_content += page_text
                    # Add page break for better text separation
                    if page_num < page_count - 1:
                        text_content += "\n\n"
                except Exception as page_error:
                    console_log(f"Error processing page {page_num}: {page_error}")
                    continue

            pdf_document.close()
            console_log(f"Successfully extracted {len(text_content)} characters from PDF: {source_name}")
            return text_content

        except Exception as e:
            console_log(f"Error extracting text from PDF {source_name}: {str(e)}")
            return f"Error extracting text from PDF: {str(e)}"

    def extract_text_from_docx(self, file_data: bytes, source_name: str) -> str:
        """Extract text from DOCX using python-docx."""
        if not DocxDocument:
            return "python-docx not available for DOCX processing"

        try:
            # Create a file-like object from bytes
            doc = DocxDocument(io.BytesIO(file_data))
            text_content = ""

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content += cell.text + " "
                    text_content += "\n"

            console_log(f"Successfully extracted text from DOCX: {source_name}")
            return text_content.strip()

        except Exception as e:
            console_log(f"Error extracting text from DOCX {source_name}: {str(e)}")
            return f"Error extracting text from DOCX: {str(e)}"

    def extract_text_from_doc(self, file_data: bytes, source_name: str) -> str:
        """Extract text from DOC (legacy format) using docx2txt."""
        if not docx2txt:
            return "docx2txt not available for DOC processing"

        try:
            # Create a temporary file for docx2txt
            with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as temp_file:
                temp_file.write(file_data)
                temp_file_path = temp_file.name

            # Extract text using docx2txt
            text_content = docx2txt.process(temp_file_path)

            # Clean up temporary file
            os.unlink(temp_file_path)

            console_log(f"Successfully extracted text from DOC: {source_name}")
            return text_content.strip()

        except Exception as e:
            console_log(f"Error extracting text from DOC {source_name}: {str(e)}")
            return f"Error extracting text from DOC: {str(e)}"

    def extract_text_from_pptx(self, file_data: bytes, source_name: str) -> str:
        """Extract text from PPTX using python-pptx."""
        if not Presentation:
            return "python-pptx not available for PPTX processing"

        try:
            # Create a file-like object from bytes
            prs = Presentation(io.BytesIO(file_data))
            text_content = ""

            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and getattr(shape, "text", "").strip():
                        slide_text += getattr(shape, "text", "") + " "

                # Extract text from tables
                for shape in slide.shapes:
                    if hasattr(shape, "has_table") and getattr(shape, "has_table", False):
                        table = getattr(shape, "table", None)
                        if table and hasattr(table, "rows"):
                            for row in table.rows:
                                if hasattr(row, "cells"):
                                    for cell in row.cells:
                                        if hasattr(cell, "text") and getattr(cell, "text", "").strip():
                                            slide_text += getattr(cell, "text", "") + " "

                if slide_text.strip():
                    text_content += f"Slide {slide_num + 1}:\n{slide_text.strip()}\n\n"

            console_log(f"Successfully extracted text from PPTX: {source_name}")
            return text_content.strip()

        except Exception as e:
            console_log(f"Error extracting text from PPTX {source_name}: {str(e)}")
            return f"Error extracting text from PPTX: {str(e)}"

    def extract_text_from_ppt(self, file_data: bytes, source_name: str) -> str:
        """Extract text from PPT (legacy format) using python-pptx.
        Note: This is a fallback as python-pptx primarily supports PPTX."""
        try:
            # Try to process as PPTX first
            return self.extract_text_from_pptx(file_data, source_name)
        except Exception as e:
            console_log(f"Error extracting text from PPT {source_name}: {str(e)}")
            return f"Error extracting text from PPT: {str(e)}"

    def extract_text_from_json(self, file_data: bytes, source_name: str) -> str:
        """Extract text from JSON by formatting it as readable text."""
        try:
            # Decode bytes to string
            json_str = file_data.decode("utf-8")

            # Parse JSON
            json_data = json.loads(json_str)

            # Convert to formatted string
            formatted_text = json.dumps(json_data, indent=2, ensure_ascii=False)

            console_log(f"Successfully extracted text from JSON: {source_name}")
            return formatted_text

        except Exception as e:
            console_log(f"Error extracting text from JSON {source_name}: {str(e)}")
            return f"Error extracting text from JSON: {str(e)}"

    def extract_text_from_file(self, file_data: bytes, file_type: str, source_name: str) -> str:
        """Extract text from file based on file type."""
        file_type = file_type.lower()

        if file_type == "pdf":
            return self.extract_text_from_pdf(file_data, source_name)
        elif file_type == "docx":
            return self.extract_text_from_docx(file_data, source_name)
        elif file_type == "doc":
            return self.extract_text_from_doc(file_data, source_name)
        elif file_type == "pptx":
            return self.extract_text_from_pptx(file_data, source_name)
        elif file_type == "ppt":
            return self.extract_text_from_ppt(file_data, source_name)
        elif file_type == "json":
            return self.extract_text_from_json(file_data, source_name)
        else:
            console_log(f"Unsupported file type: {file_type}")
            return f"Unsupported file type: {file_type}"

    def get_supported_formats(self) -> list:
        """Get list of supported file formats."""
        return ["pdf", "docx", "doc", "pptx", "ppt", "json"]
